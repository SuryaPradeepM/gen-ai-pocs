import subprocess
from pptx import Presentation
from PyPDF2 import PdfReader
from langchain_groq import ChatGroq
from langchain_ollama.llms import OllamaLLM
from langchain.chat_models import AzureChatOpenAI
from langchain.prompts import (
    ChatPromptTemplate,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
    MessagesPlaceholder,
)
from langchain.chains import ConversationChain
from langchain.memory import ConversationBufferMemory
import os
from dotenv import load_dotenv
import asyncio
import azure.cognitiveservices.speech as speechsdk


load_dotenv()

# TODO make system msg
# TODO install ollama and all

def convert_pptx_to_pdf(pptx_path):
    """
    Converts a PowerPoint file to PDF using comtypes (Windows) or pdf2pptx (other platforms).
    """
    pdf_path = pptx_path.replace(".pptx", ".pdf")

    try:
        if os.name == 'nt':  # Windows
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            powerpoint.Visible = True
            deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
            deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 represents PDF format
            deck.Close()
            powerpoint.Quit()
        else:  # Other platforms
            # TODO: To test
            from pdf2pptx.tools import convert
            convert(pptx_path, pdf_path)

        return pdf_path

    except ImportError as e:
        raise ImportError(
            "Required conversion package not found. Please install comtypes (Windows) or pdf2pptx (other platforms)"
        ) from e
    except Exception as e:
        raise RuntimeError(f"Failed to convert PPTX to PDF: {str(e)}") from e


class PPTSlideIterator:
    _instances = {}

    def __new__(cls, file_path):
        if file_path not in cls._instances:
            instance = super().__new__(cls)
            instance.pptx_path = file_path
            instance.pdf_path = convert_pptx_to_pdf(file_path)
            # instance.pdf_path = ""
            instance.presentation = Presentation(file_path)
            instance.pdf_reader = PdfReader(instance.pdf_path)
            instance.current_slide = 1
            instance.total_slides = len(instance.presentation.slides)
            cls._instances[file_path] = instance
        return cls._instances[file_path]

    def __init__(self, file_path):
        # Initialize is called after __new__, but we don't need to do anything here
        # since initialization is handled in __new__
        pass

    def get_slide_content(self, slide_number):
        if 1 <= slide_number <= self.total_slides:
            slide = self.presentation.slides[slide_number - 1]
            content = f"Slide {slide_number}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += f"{shape.text}\n"
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text
                if notes_text.strip():
                    content += f"Notes:\n{notes_text}\n"
            return content
        else:
            return None

    def get_pdf_page(self, slide_number):
        if 1 <= slide_number <= self.total_slides:
            return self.pdf_reader.pages[slide_number - 1]
            # return 11
        else:
            return None

    def get_total_slides(self):
        return self.total_slides

    def get_slides_content_up_to(self, slide_number):
        content = ""
        for i in range(1, slide_number + 1):
            content += self.get_slide_content(i) + "\n\n"
        return content.strip()

    def get_entire_presentation_content(self):
        return self.get_slides_content_up_to(self.total_slides)


def load_llm(model):
    if model == "groq":
        GROQ_API_KEY = os.getenv("GROQ_API_KEY")
        return ChatGroq(
            temperature=0,
            model="llama3-70b-8192",
            groq_api_key=GROQ_API_KEY,
            max_tokens=50,
        )
    elif model == "ollama-llama8b":
        return OllamaLLM(model=model)
    elif model == "gpt-4":
        return AzureChatOpenAI(
            openai_api_version=os.getenv("AZURE_OPENAI_API_VERSION"),
            azure_deployment=os.getenv("AZURE_OPENAI_GPT_DEPLOYMENT_NAME"),
            temperature=0,
        )
    elif model in ["llama3:latest", "phi3:mini-128k"]:
        return OllamaLLM(model=model)
    elif model == "gpt-4-o":
        os.environ["OPENAI_API_VERSION"] = os.environ["AZURE_OPENAI_API_VERSION"]
        os.environ["AZURE_OPENAI_ENDPOINT"] = os.environ["AZURE_OPENAI_ENDPOINT"]
        os.environ["AZURE_OPENAI_API_KEY"] = os.environ["AZURE_OPENAI_API_KEY"]
        return AzureChatOpenAI(
            deployment_name=os.environ["AZURE_OPENAI_GPT_DEPLOYMENT_NAME"]
        )
    else:
        raise ValueError(f"Unsupported model: {model}")


async def explain_slide(pp, slide_number, llm):
    slide_content = pp.get_slide_content(slide_number)
    system_template = """You are an AI assistant explaining a PowerPoint presentation. 
    Explain the following slide content in a clear and concise manner. 
    After your explanation, ask if the user has any questions about this slide.
    Don't start your presentation with this slide. Just present the slide creatively and stick to the context in the slide."""

    human_template = "{slide_content}"

    chat_prompt = ChatPromptTemplate.from_messages(
        [
            SystemMessagePromptTemplate.from_template(system_template),
            HumanMessagePromptTemplate.from_template(human_template),
        ]
    )

    messages = chat_prompt.format_prompt(slide_content=slide_content).to_messages()
    response = await asyncio.to_thread(llm.invoke, messages)

    try:
        return response.content
    except:
        return response


def initiate_bot(pp, llm, convo_history):
    system_prompt = "Role: You are a conversational agent acting as an expert in presenting PowerPoint slides and answering questions about the presentation content. Your name is Saroja."

    docs = pp.get_entire_presentation_content()

    human_template = f"""Task: You are provided with information present in the PowerPoint document: 

    {docs}

    Your task is to answer questions asked by the user related to any of these slides. Keep conversations shorter and on point. Note strictly that you are a two-way conversational agent."""

    chat_prompt = ChatPromptTemplate.from_messages(
        [
            SystemMessagePromptTemplate.from_template(system_prompt),
            MessagesPlaceholder(variable_name="history"),
            HumanMessagePromptTemplate.from_template(
                human_template + "\n\nHuman: {input}\nAI:"
            ),
        ]
    )

    memory = ConversationBufferMemory(k=6, return_messages=True)

    for entry in convo_history:
        memory.save_context({"input": entry["human"]}, {"output": entry["ai"]})

    conversation = ConversationChain(
        llm=llm, verbose=True, prompt=chat_prompt, memory=memory
    )

    return conversation, convo_history


async def chat(conversation, user_input, convo_history):
    response = await asyncio.to_thread(conversation.predict, input=user_input)
    convo_history.append({"human": user_input, "ai": response})
    return response, convo_history


async def azure_text_to_speech(text: str) -> str:
    speech_config = speechsdk.SpeechConfig(
        subscription=os.getenv("AZURE_SPEECH_KEY"),
        region=os.getenv("AZURE_SPEECH_REGION"),
    )

    if not speech_config:
        raise Exception(
            "Speech configuration could not be initialized. Check your Azure credentials."
        )

    audio_config = speechsdk.audio.AudioOutputConfig(filename="response.mp3")
    speech_synthesizer = speechsdk.SpeechSynthesizer(
        speech_config=speech_config, audio_config=audio_config
    )

    # Start the synthesis and wait for it to complete
    loop = asyncio.get_event_loop()
    future = loop.run_in_executor(
        None, lambda: speech_synthesizer.speak_text_async(text).get()
    )
    result = await future

    if result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
        return "response.mp3"
    else:
        raise Exception(
            "Speech synthesis failed: Reason: {} Error Details: {}".format(
                result.cancellation_details.reason,
                result.cancellation_details.error_details
                if result.cancellation_details
                else "None",
            )
        )


if __name__ == "__main__":
    pp = PPTSlideIterator("./Prompt Engineering 101 -v1.pptx")
    print(pp)
    print(pp.pptx_path)
    print(pp.pdf_path)
