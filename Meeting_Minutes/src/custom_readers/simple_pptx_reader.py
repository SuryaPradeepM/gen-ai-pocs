"""Slides parser.

Contains Simple Parser for .pptx files.

Source: https://github.com/run-llama/llama_index/blob/main/llama-index-integrations/readers/llama-index-readers-file/llama_index/readers/file/slides/base.py

Changes made to remove logic which does image captioning
(Firewall issues with huggingface models & additional latency without significant improvement)
"""

from pathlib import Path
from typing import Dict, List, Optional

from fsspec import AbstractFileSystem
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from pptx import Presentation
from tqdm import tqdm


class SimplePptxReader(BaseReader):
    """Simpler Powerpoint parser.

    Extracts text, *does not caption images*, and specify slides.
    """

    def __init__(self, input_files: Optional[List] = None) -> None:
        """Init Simple PPT parser.
        Args:
            input_files (Optional[List], optional): List of file paths to read. Defaults to None.
        """
        self.input_files = input_files

    def load_data(
        self,
        file: Path,
        extra_info: Optional[Dict] = None,
        fs: Optional[AbstractFileSystem] = None,
    ) -> List[Document]:
        """Loads the file and parses them to return document for each input_file

        Args:
            file (Path): Path to input file
            extra_info (Optional[Dict], optional): A dictionary containing extra metadata information. Defaults to None.
            fs (Optional[fsspec.AbstractFileSystem]): File system to use. Defaults
                to using the local file system. Can be changed to use any remote file system
                exposed via the fsspec interface.
        Returns:
            List[Document]: A list of documents loaded from the input files.
        """
        if fs:
            with fs.open(file) as f:
                presentation = Presentation(f)
        else:
            presentation = Presentation(file)

        try:
            metadata = {"file_name": file.name}
        except:
            metadata = {}
        if extra_info is not None:
            metadata.update(extra_info)

        result = ""
        for i, slide in enumerate(presentation.slides):
            result += f"\n\nSlide #{i}: \n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    result += f"{shape.text}\n"
        return [Document(text=result, metadata=metadata)]

    def load_files(
        self,
        show_progress: bool = False,
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """Loads the ppt_files and parses them to return document for each input_file

        Args:
            show_progress (bool): Whether to show tqdm progress bars. Defaults to False.
            extra_info (Optional[Dict], optional): A dictionary containing extra information. Defaults to None.

        Returns:
            List[Document]: A list of documents loaded from the input files.
        """
        documents = []
        files_to_process = self.input_files
        if show_progress:
            files_to_process = tqdm(
                self.input_files, desc="Loading PPT files", unit="file"
            )
        for input_file in files_to_process:
            documents.extend(self.load_data(input_file, extra_info=extra_info))
        return documents
